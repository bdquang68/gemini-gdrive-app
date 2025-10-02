# file: upload_app.py
import os, tempfile, zipfile
import streamlit as st
import google.generativeai as genai

from langchain_text_splitters import RecursiveCharacterTextSplitter
from PyPDF2 import PdfReader
import docx, openpyxl
from pptx import Presentation

# ========= 1) Cấu hình Gemini =========
API_KEY = st.secrets.get("GOOGLE_API_KEY") or os.getenv("GOOGLE_API_KEY")
if not API_KEY:
    st.error("Thiếu GOOGLE_API_KEY. Vào App settings → Secrets để thêm.")
    st.stop()
genai.configure(api_key=API_KEY)

# ========= 2) Các hàm đọc file =========
def read_pdf(p):
    try:
        txt=""; r=PdfReader(p)
        for page in r.pages: txt += page.extract_text() or ""
        return txt
    except Exception as e: return f"\n[PDF lỗi {os.path.basename(p)}: {e}]\n"

def read_docx(p):
    try:
        d=docx.Document(p)
        return "\n".join((para.text or "") for para in d.paragraphs)
    except Exception as e: return f"\n[DOCX lỗi {os.path.basename(p)}: {e}]\n"

def read_excel(p):
    try:
        wb=openpyxl.load_workbook(p, data_only=True); out=[]
        for s in wb.sheetnames:
            for row in wb[s].iter_rows(values_only=True):
                cells=[str(c) for c in row if c is not None]
                if cells: out.append(" | ".join(cells))
        return "\n".join(out)
    except Exception as e: return f"\n[XLSX lỗi {os.path.basename(p)}: {e}]\n"

def read_pptx(p):
    try:
        prs=Presentation(p); out=[]
        for slide in prs.slides:
            for sh in slide.shapes:
                if hasattr(sh,"text"): out.append(sh.text or "")
        return "\n".join(out)
    except Exception as e: return f"\n[PPTX lỗi {os.path.basename(p)}: {e}]\n"

def read_txt(p):
    try: return open(p, "r", encoding="utf-8", errors="ignore").read()
    except Exception as e: return f"\n[TXT lỗi {os.path.basename(p)}: {e}]\n"

def read_file(p):
    pl=p.lower()
    if pl.endswith(".pdf"):  return read_pdf(p)
    if pl.endswith(".docx"): return read_docx(p)
    if pl.endswith(".xlsx"): return read_excel(p)
    if pl.endswith(".pptx"): return read_pptx(p)
    if pl.endswith(".txt") or pl.endswith(".csv"): return read_txt(p)
    return ""

# ========= 3) Chunk =========
def chunk_text(text, chunk_size=1000, overlap=200):
    return RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, chunk_overlap=overlap
    ).split_text(text)

# ========= 4) UI =========
st.title("📤 Phân tích dữ liệu từ MÁY TÍNH (Upload files hoặc .zip)")

tab_files, tab_zip = st.tabs(["Upload nhiều file", "Upload thư mục .zip"])

data = ""
with tab_files:
    uploads = st.file_uploader(
        "Chọn nhiều file (pdf, docx, xlsx, pptx, txt, csv)",
        type=["pdf","docx","xlsx","pptx","txt","csv"],
        accept_multiple_files=True
    )
    if uploads:
        with tempfile.TemporaryDirectory() as tmp:
            texts=[]
            for uf in uploads:
                save_path = os.path.join(tmp, uf.name)
                with open(save_path, "wb") as f: f.write(uf.read())
                texts.append(read_file(save_path))
            data = "\n".join(texts)
        st.success(f"✅ Đã nạp {len(uploads)} file.")

with tab_zip:
    zip_file = st.file_uploader("Upload 1 file .zip chứa cả thư mục", type=["zip"])
    if zip_file:
        with tempfile.TemporaryDirectory() as tmp:
            zpath = os.path.join(tmp, "in.zip")
            with open(zpath, "wb") as f: f.write(zip_file.getvalue())
            with zipfile.ZipFile(zpath, 'r') as zf:
                zf.extractall(tmp)
            texts=[]
            for root, _, files in os.walk(tmp):
                for name in files:
                    p = os.path.join(root, name)
                    if p.endswith("in.zip"): continue
                    texts.append(read_file(p))
            data = "\n".join(texts)
        st.success("✅ Đã giải nén & nạp dữ liệu từ .zip.")
# ---- Làm chữ trong ô nhập đen và đậm hơn ----
st.markdown("""
<style>
/* Ô nhập text */
.stTextInput input {
    color: #000000 !important;        /* chữ đen hẳn */
    font-weight: 700 !important;      /* đậm */
    font-size: 16px !important;       /* to hơn chút (tùy chỉnh) */
}

/* Placeholder trong ô nhập */
.stTextInput input::placeholder {
    color: #555555 !important;        /* placeholder cũng đậm hơn mặc định */
    font-weight: 500 !important;
}
</style>
""", unsafe_allow_html=True)

# ========= 5) Hỏi AI =========
query = st.text_input("Nhập câu hỏi:")
if query and data:
    try:
        model = genai.GenerativeModel(
            "gemini-1.5-pro",
            generation_config={"temperature": 0.3, "max_output_tokens": 1024}
        )

        # Thu gọn context để tránh 400 "prompt too long"
        context = data[:20000]
        prompt = f"Dữ liệu:\n{context}\n\nCâu hỏi: {query}\n\nTrả lời ngắn gọn, dựa vào dữ liệu."

        with st.spinner("Đang phân tích bằng Gemini..."):
            resp = model.generate_content(prompt)

        # Một số case bị chặn nội dung → resp.text trống
        answer = getattr(resp, "text", "") or ""
        if not answer:
            # Thử đọc candidate/errors để show rõ lý do
            blocked = []
            for c in getattr(resp, "candidates", []) or []:
                safety = getattr(c, "safety_ratings", None) or []
                if safety:
                    blocked.extend([f"{r.category}={r.probability}" for r in safety])
            if blocked:
                st.error("Nội dung bị chặn bởi bộ lọc an toàn của Gemini:\n" + "\n".join(blocked))
            else:
                st.warning("Gemini không trả lời (resp.text rỗng). Thử rút gọn dữ liệu hoặc đổi câu hỏi.")
        else:
            st.subheader("🔎 Kết quả AI phân tích")
            st.write(answer)

    except genai.types.generation_types.BlockedPromptException as e:
        st.exception(e)
        st.error("Prompt bị chặn bởi bộ lọc an toàn. Hãy rút gọn hoặc điều chỉnh nội dung.")
    except genai.types.generation_types.StopCandidateException as e:
        st.exception(e)
        st.error("Gemini dừng sinh nội dung sớm. Thử hỏi ngắn hơn.")
    except Exception as e:
        # Hiện toàn bộ lỗi để dễ khoanh vùng (HTTP 400/401/403/429/500…)
        st.exception(e)
        st.error(
            "Lỗi khi gọi Gemini. Có thể do:\n"
            "• GOOGLE_API_KEY sai/hết hạn/quota (401/403/429)\n"
            "• Prompt quá dài (400) → rút gọn dữ liệu\n"
            "• Sự cố tạm thời phía server (500)\n"
        )
